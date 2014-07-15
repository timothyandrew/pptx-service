import pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import re
from io import BytesIO

class Presentation(object):
  def __init__(self, title, body):
    self.title = title
    self.body = body
    self.presentation = pptx.Presentation()

  def add_slide(self, text, bold = False):
    slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(-1), Inches(15), Inches(10))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    left = top = width = height = Inches(1)
    content = slide.shapes.add_textbox(Inches(4.5), Inches(2), width, height)

    textframe = content.textframe
    textframe.clear()
    textframe.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    textframe.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = textframe.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    font = run.font
    font.name = 'Constantia'
    font.size = Pt(40)
    font.color.rgb = RGBColor(255,255,255)
    font.bold = bold

    run.text = text

  def generate(self):
    self.add_slide(self.title, True)
    for text in re.split("\r\n\r\n", self.body):
      text = text.replace("\r\n", "\n")
      self.add_slide(text)

    buf = BytesIO()
    self.presentation.save(buf)
    buf.seek(0)
    return buf