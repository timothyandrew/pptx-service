from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import re

prs = Presentation()

def add_title_slide(prs, text):
  slide = prs.slides.add_slide(prs.slide_layouts[0])
  left = top = width = height = Inches(1)
  content = slide.shapes[0]

  textframe = content.textframe
  textframe.clear()
  textframe.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
  textframe.vertical_anchor = MSO_ANCHOR.TOP
  paragraph = textframe.paragraphs[0]
  paragraph.alignment = PP_ALIGN.CENTER
  run = paragraph.add_run()
  font = run.font
  font.name = 'Constantia'
  font.size = Pt(40)
  font.bold = True

  run.text = text

def add_slide(prs, text):
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  left = top = width = height = Inches(1)
  content = slide.shapes.add_textbox(Inches(4.5), Inches(2), width, height)

  textframe = content.textframe
  textframe.clear()
  textframe.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
  textframe.vertical_anchor = MSO_ANCHOR.TOP
  paragraph = textframe.paragraphs[0]
  paragraph.alignment = PP_ALIGN.CENTER
  run = paragraph.add_run()
  font = run.font
  font.name = 'Constantia'
  font.size = Pt(40)

  run.text = text

text = """First

Second

Third
"""

add_title_slide(prs, "Title Slide")
for text in re.split("\n\n", text):
  add_slide(prs, text)


prs.save('/tmp/test.pptx')